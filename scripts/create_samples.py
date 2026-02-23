import sys
import subprocess

def install(package):
    subprocess.check_call([sys.executable, "-m", "pip", "install", package])

try:
    from pptx import Presentation
except ImportError:
    install("python-pptx")
    from pptx import Presentation

try:
    import openpyxl
except ImportError:
    install("openpyxl")
    import openpyxl

from pathlib import Path

data_dir = Path("/home/nobuhiko/project/Docling_pdf2markdown/tests/test_data")
data_dir.mkdir(parents=True, exist_ok=True)

# Generate PPTX
pptx_path = data_dir / "sample.pptx"
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Docling Test Presentation"
subtitle.text = "Testing native PPTX support in Docling v2"

# Add a slide with a table
blank_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(blank_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
title_shape.text = "Sample Table"

rows = 3
cols = 2
left = top = width = height = 100
table = shapes.add_table(rows, cols, left, top, width, height).table
table.columns[0].width = 1000000
table.columns[1].width = 1000000

table.cell(0, 0).text = "Project"
table.cell(0, 1).text = "Status"
table.cell(1, 0).text = "Docling V2"
table.cell(1, 1).text = "Done"
table.cell(2, 0).text = "Testing"
table.cell(2, 1).text = "In Progress"

prs.save(pptx_path)
print(f"Created {pptx_path}")

# Generate XLSX
xlsx_path = data_dir / "sample.xlsx"
wb = openpyxl.Workbook()
ws = wb.active
ws.title = "Test Data"
ws["A1"] = "ID"
ws["B1"] = "Name"
ws["C1"] = "Score"

data = [
    (1, "Alice", 95),
    (2, "Bob", 88),
    (3, "Charlie", 92),
]

for row in data:
    ws.append(row)

wb.save(xlsx_path)
print(f"Created {xlsx_path}")
